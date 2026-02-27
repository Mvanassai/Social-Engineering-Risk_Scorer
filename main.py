import io
import re
import uvicorn
import spacy
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

# --- MODEL LOADING ---
try:
    nlp = spacy.load("en_core_web_lg")
except Exception:
    print("Warning: spaCy model 'en_core_web_lg' not found. Using Regex Fallback.")
    nlp = None

# --- FILE SUPPORT LIBRARIES ---
try:
    import fitz  # PyMuPDF for PDF
except ImportError:
    fitz = None
    print("Warning: PyMuPDF (fitz) not installed. PDF support disabled.")

try:
    from docx import Document  # python-docx for .docx
except ImportError:
    Document = None
    print("Warning: python-docx not installed. DOCX support disabled.")

try:
    from pptx import Presentation  # python-pptx for .ppt / .pptx
except ImportError:
    Presentation = None
    print("Warning: python-pptx not installed. PPT/PPTX support disabled. Install with: pip install python-pptx")

app = FastAPI(title="Sentinel OSINT - Privacy Hardening API")

# --- CORS ---
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- CORE LOGIC ENGINE (unchanged) ---
class SentinelCore:
    def __init__(self):
        self.taxonomy = {
            "google": "a Tier-1 Tech Corporation",
            "microsoft": "a Global Software Leader",
            "manager": "Strategic Director",
            "engineer": "Technical Specialist",
            "sql": "Structured Database Systems",
            "hyderabad": "a major Tech Hub"
        }

    def simulate_adversary(self, text: str):
        logs = []
        score = 0
        
        if re.search(r'[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+', text):
            score += 40
            logs.append("LEAK: Direct email identified. Vulnerable to Phishing.")
        
        if re.search(r'(\+?\d{1,3}[\s-]?)?(\d{10})', text):
            score += 30
            logs.append("LEAK: Phone contact exposed. Vulnerable to SIM-swap.")
            
        for brand in self.taxonomy.keys():
            if brand in text.lower():
                score += 15
                logs.append(f"RECON: Connection to {brand.upper()} found.")
                
        return min(score, 100), logs

    def nlp_harden(self, text: str):
        if nlp is None:
            return self.fallback_harden(text)
        
        doc = nlp(text)
        hardened = text
        for ent in reversed(doc.ents):
            if ent.label_ in ["PERSON", "ORG", "GPE", "FAC", "LOC"]:
                replacement = f"[HIDDEN_{ent.label_}]"
                hardened = hardened[:ent.start_char] + replacement + hardened[ent.end_char:]
        
        hardened = re.sub(r'[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+', "[ID_GATEWAY]", hardened)
        hardened = re.sub(r'(\+?\d{1,3}[\s-]?)?(\d{10})', "[VERIFIED_LINE]", hardened)
        return hardened

    def fallback_harden(self, text: str):
        safe = text
        safe = re.sub(r'[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+', "[ID_GATEWAY]", safe)
        safe = re.sub(r'(\+?\d{1,3}[\s-]?)?(\d{10})', "[VERIFIED_LINE]", safe)
        for k, v in self.taxonomy.items():
            pattern = re.compile(re.escape(k), re.IGNORECASE)
            safe = pattern.sub(v, safe)
        return safe

# --- DATA MODELS ---
class DataIn(BaseModel):
    content: str

class UrlIn(BaseModel):
    url: str

# --- API ENDPOINTS ---

@app.post("/analyze-url")
async def analyze_url(data: UrlIn):
    url = data.url.lower()
    engine = SentinelCore()
    
    if "linkedin.com" in url:
        username = url.split("/in/")[-1].strip("/")
        sim_text = f"LinkedIn Profile: {username}. Manager at Google and Microsoft. Reach me at {username}@google.com"
        orig_risk = 85
    elif "twitter.com" in url or "x.com" in url:
        handle = url.split("/")[-1].strip("@")
        sim_text = f"X handle: @{handle}. Senior Engineer at Microsoft. Email: {handle}@outlook.com"
        orig_risk = 75
    elif "facebook.com" in url:
        user = url.split("/")[-1] if "/" in url else "User"
        sim_text = f"Facebook: {user}. Lives in Hyderabad. Works at Google. Contact: 9876543210"
        orig_risk = 95
    else:
        raise HTTPException(status_code=400, detail="URL not supported. Use LinkedIn, X, or Facebook.")

    risk, logs = engine.simulate_adversary(sim_text)
    
    return {
        "extracted_text": sim_text,
        "risk_score": risk,
        "original_risk": orig_risk,
        "market_score": 100 - risk,
        "evidence": logs,
        "safe_text": engine.nlp_harden(sim_text)
    }

@app.post("/process")
async def process(data: DataIn):
    engine = SentinelCore()
    risk, logs = engine.simulate_adversary(data.content)
    
    return {
        "original_risk": risk,
        "risk_score": risk,
        "market_score": 100 - risk,
        "evidence": logs,
        "safe_text": engine.nlp_harden(data.content)
    }

@app.post("/upload")
async def upload(file: UploadFile = File(...)):
    """Extracts text from PDF, DOCX, PPT/PPTX and runs analysis"""
    body = await file.read()
    text = ""
    filename = file.filename.lower()

    try:
        if filename.endswith(".docx") and Document is not None:
            doc = Document(io.BytesIO(body))
            text = "\n".join([p.text for p in doc.paragraphs])

        elif filename.endswith(".pdf") and fitz is not None:
            pdf = fitz.open(stream=body, filetype="pdf")
            text = "".join([page.get_text() for page in pdf])

        elif filename.endswith((".ppt", ".pptx")) and Presentation is not None:
            prs = Presentation(io.BytesIO(body))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text += run.text + " "
                    text += "\n"

        else:
            # Fallback for plain text or unsupported
            text = body.decode("utf-8", errors="ignore")

    except Exception as e:
        raise HTTPException(status_code=400, detail=f"File processing error: {str(e)}")

    if not text.strip():
        raise HTTPException(status_code=400, detail="No readable text found in the file.")

    engine = SentinelCore()
    risk, logs = engine.simulate_adversary(text)

    return {
        "extracted_text": text[:1500],  # Increased preview length
        "original_risk": risk,
        "risk_score": risk,
        "market_score": 100 - risk,
        "evidence": logs,
        "safe_text": engine.nlp_harden(text)
    }

if __name__ == "__main__":
    uvicorn.run(app, host="127.0.0.1", port=8000)